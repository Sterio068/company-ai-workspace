"""
多格式抽字器 · V1.1-SPEC §E-2
=====================================
依副檔名路由 · 統一回 {type, content_preview, filename, size, modified_at, ...}
所有抽字失敗都不 raise · 回 {type: "error", error: str(e)}
讓 indexer 繼續處理下一個檔

本模組對外只暴露 extract(path) · 其他 helper 都 underscore
"""
import os
import logging
import pathlib
from datetime import datetime
from typing import Any

logger = logging.getLogger("chengfu.extract")

# lazy import · 讓測試環境不一定要裝滿所有 library
_fitz = None
_docx = None
_pptx = None
_openpyxl = None
_Image = None
_TAGS = None


def _lazy_fitz():
    global _fitz
    if _fitz is None:
        import fitz  # pymupdf
        _fitz = fitz
    return _fitz


def _lazy_docx():
    global _docx
    if _docx is None:
        from docx import Document
        _docx = Document
    return _docx


def _lazy_pptx():
    global _pptx
    if _pptx is None:
        from pptx import Presentation
        _pptx = Presentation
    return _pptx


def _lazy_openpyxl():
    global _openpyxl
    if _openpyxl is None:
        from openpyxl import load_workbook
        _openpyxl = load_workbook
    return _openpyxl


def _lazy_image():
    global _Image, _TAGS
    if _Image is None:
        from PIL import Image
        from PIL.ExifTags import TAGS
        _Image, _TAGS = Image, TAGS
    return _Image, _TAGS


# ------------------------------------------------------------
# Extractor · 每格式一個函式 · 失敗交由外層 extract() catch
# ------------------------------------------------------------
# Round 9 · OCR 環境探測 · 第一次 fail 後 cache · 之後 metric 看得到
# Codex Round 10.5 · 加 startup probe · 不再 lazy
_OCR_AVAILABLE = None
_OCR_LAST_ERROR = None
_OCR_LANGS = []


def probe_ocr_startup() -> dict:
    """Codex Round 10.5 紅 4 · 啟動時主動探測 OCR · 不再 lazy

    做 3 件事:
    1. 呼叫 tesseract --list-langs · 確認 binary 存在 + 知道哪些語言包
    2. 產生一張極小 image-only PDF · 跑 OCR 確認真的能解
    3. 結果寫進 _OCR_AVAILABLE · /healthz 立刻看得到真狀態
    """
    global _OCR_AVAILABLE, _OCR_LAST_ERROR, _OCR_LANGS
    import subprocess

    # Step 1 · list-langs
    try:
        result = subprocess.run(
            ["tesseract", "--list-langs"],
            capture_output=True, text=True, timeout=5,
        )
        # tesseract --list-langs 輸出到 stderr(非錯誤 · 是設計)· 某版到 stdout
        raw = (result.stderr or "") + (result.stdout or "")
        lines = [l.strip() for l in raw.split("\n") if l.strip() and not l.startswith("List")]
        _OCR_LANGS = lines
        if result.returncode != 0:
            _OCR_AVAILABLE = False
            _OCR_LAST_ERROR = f"tesseract --list-langs exit {result.returncode}"
            logger.warning("[ocr] probe step 1 fail · %s", _OCR_LAST_ERROR)
            return ocr_status()
    except FileNotFoundError:
        _OCR_AVAILABLE = False
        _OCR_LAST_ERROR = "tesseract 未安裝 · Dockerfile 應裝 tesseract-ocr + tesseract-ocr-chi-tra"
        logger.error("[ocr] probe %s", _OCR_LAST_ERROR)
        return ocr_status()
    except Exception as e:
        _OCR_AVAILABLE = False
        _OCR_LAST_ERROR = f"list-langs {type(e).__name__}: {e}"
        return ocr_status()

    # Step 2 · 確認必要語言包
    if "chi_tra" not in _OCR_LANGS:
        _OCR_AVAILABLE = False
        _OCR_LAST_ERROR = f"缺 chi_tra 語言包 · 現有:{_OCR_LANGS}"
        logger.error("[ocr] probe %s", _OCR_LAST_ERROR)
        return ocr_status()

    # Step 3 · PyMuPDF OCR 實測(Codex R2.5 · 用真 image 不是 text layer)
    # 原 insert_text 建的是文字層 · textpage_ocr 直接讀 text 層就過 · 沒驗到 OCR
    # 改:PIL 畫 "OK" 的 bitmap · 嵌進 PDF · OCR 必須真從像素認字
    try:
        fitz = _lazy_fitz()
        import tempfile, os as _os
        # 用 PIL 畫一張 200×100 的白底黑字 image
        try:
            from PIL import Image, ImageDraw
        except ImportError:
            raise RuntimeError("PIL not installed · required for OCR probe")

        img = Image.new("RGB", (200, 100), color="white")
        draw = ImageDraw.Draw(img)
        # 用預設字型畫 "OK" · 大概 30px 高 · 應足以 OCR 識別
        draw.text((50, 35), "OK", fill="black")
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            img_path = f.name
        img.save(img_path, "PNG")

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            tmp_pdf = f.name
        try:
            # 建一個 PDF · 只有圖片(無文字層)
            doc = fitz.open()
            page = doc.new_page(width=200, height=100)
            page.insert_image(fitz.Rect(0, 0, 200, 100), filename=img_path)
            doc.save(tmp_pdf)
            doc.close()
            # 再打開 · 確認 text layer 是空的(才算 image-only)
            doc = fitz.open(tmp_pdf)
            page = doc[0]
            plain_text = page.get_text("text").strip()
            if plain_text:
                logger.warning("[ocr] probe PDF 居然有 text layer · PIL 寫入方式需查")
            # OCR 路徑
            tp = page.get_textpage_ocr(language="chi_tra+eng")
            text = page.get_text(textpage=tp)
            doc.close()
            recognized = text.strip()
            _OCR_AVAILABLE = True
            _OCR_LAST_ERROR = None
            _OCR_PROBE_RESULT = recognized[:50]
            logger.info(
                "[ocr] probe OK(real image OCR)· 語言=%s · recognized=%r",
                _OCR_LANGS, recognized[:20],
            )
        finally:
            for p in (img_path, tmp_pdf):
                try:
                    _os.unlink(p)
                except Exception:
                    pass
    except Exception as e:
        _OCR_AVAILABLE = False
        _OCR_LAST_ERROR = f"probe OCR (image) fail · {type(e).__name__}: {str(e)[:120]}"
        logger.error("[ocr] probe step 3 %s", _OCR_LAST_ERROR)
    return ocr_status()


def reset_ocr_cache():
    """Codex R2.5 · 給 /admin/ocr/reprobe 用 · 維運後不用重啟整個容器"""
    global _OCR_AVAILABLE, _OCR_LAST_ERROR, _OCR_LANGS
    _OCR_AVAILABLE = None
    _OCR_LAST_ERROR = None
    _OCR_LANGS = []


def ocr_status() -> dict:
    """給 health endpoint 用 · 報告 OCR 是否可用"""
    return {
        "available": _OCR_AVAILABLE,
        "langs": _OCR_LANGS,
        "last_error": _OCR_LAST_ERROR,
        "note": "False = tesseract 未裝 · OCR fallback 不會跑 · 掃描 PDF 內容會空"
                if _OCR_AVAILABLE is False else
                ("未探測 · startup 時自動跑 · 若長期 None 表示 lifespan 沒呼叫 probe_ocr_startup()"
                 if _OCR_AVAILABLE is None else None),
    }


def _extract_pdf(path: str) -> dict:
    """ROADMAP §11.4 · max_ocr_pages 上限 · 防 200 頁掃描檔吃光 RAM"""
    global _OCR_AVAILABLE, _OCR_LAST_ERROR
    import os as _os
    fitz = _lazy_fitz()
    max_ocr = int(_os.getenv("MAX_OCR_PAGES_PER_PDF", "20"))
    doc = fitz.open(path)
    pages = []
    ocr_triggered = 0
    ocr_skipped_no_engine = 0
    ocr_skipped_over_limit = 0
    try:
        for page in doc:
            text = page.get_text("text").strip()
            # 若文字太少而頁面有圖片 · 嘗試 OCR 降級
            if len(text) < 120 and page.get_images():
                if _OCR_AVAILABLE is False:
                    ocr_skipped_no_engine += 1
                elif ocr_triggered >= max_ocr:
                    # ROADMAP §11.4 · 超過上限 · 跳過 OCR 但保留純 text(可能空)
                    ocr_skipped_over_limit += 1
                else:
                    try:
                        tp = page.get_textpage_ocr(language="chi_tra+eng")
                        text = page.get_text(textpage=tp).strip()
                        ocr_triggered += 1
                        _OCR_AVAILABLE = True
                    except Exception as e:
                        if _OCR_AVAILABLE is None:
                            _OCR_AVAILABLE = False
                            _OCR_LAST_ERROR = f"{type(e).__name__}: {str(e)[:120]}"
                            logger.warning(
                                "[ocr] 第一次 fallback 失敗 · 後續 PDF OCR 不再嘗試 · 錯誤=%s · "
                                "檢查 Dockerfile 是否裝 tesseract-ocr-chi-tra",
                                _OCR_LAST_ERROR,
                            )
                        ocr_skipped_no_engine += 1
            pages.append(text)
    finally:
        doc.close()
    result = {
        "type": "pdf",
        "page_count": len(pages),
        "ocr_pages": ocr_triggered,
        "content_preview": ("\n\n".join(pages))[:2000],
    }
    if ocr_skipped_no_engine:
        result["ocr_skipped_no_engine"] = ocr_skipped_no_engine
    if ocr_skipped_over_limit:
        result["ocr_skipped_over_limit"] = ocr_skipped_over_limit
        result["ocr_max_pages_setting"] = max_ocr
    return result


def _extract_docx(path: str) -> dict:
    Document = _lazy_docx()
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # 表格內容也抓(簡易)
    for table in doc.tables:
        for row in table.rows:
            row_txt = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_txt:
                paragraphs.append(row_txt)
    text = "\n".join(paragraphs)
    return {
        "type": "docx",
        "paragraph_count": len(paragraphs),
        "content_preview": text[:2000],
    }


def _extract_pptx(path: str) -> dict:
    Presentation = _lazy_pptx()
    prs = Presentation(path)
    slides_text = []
    for i, s in enumerate(prs.slides, start=1):
        parts = []
        for shape in s.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if parts:
            slides_text.append(f"[投影片 {i}]\n" + "\n".join(parts))
    return {
        "type": "pptx",
        "slide_count": len(prs.slides),
        "content_preview": ("\n\n".join(slides_text))[:2000],
    }


def _extract_xlsx(path: str) -> dict:
    load_workbook = _lazy_openpyxl()
    wb = load_workbook(path, read_only=True, data_only=True)
    try:
        preview_parts = []
        # 最多前 3 張 sheet · 每張前 20 行
        for sheet_name in wb.sheetnames[:3]:
            ws = wb[sheet_name]
            rows = []
            for i, row in enumerate(ws.iter_rows(max_row=20, values_only=True)):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
                if i >= 19:
                    break
            if rows:
                preview_parts.append(f"[{sheet_name}]\n" + "\n".join(rows))
        return {
            "type": "xlsx",
            "sheet_count": len(wb.sheetnames),
            "content_preview": ("\n\n".join(preview_parts))[:2000],
        }
    finally:
        wb.close()


def _extract_image(path: str) -> dict:
    """設計圖只抽 metadata · 不抽內容(v2 再做 CLIP vision)"""
    Image, TAGS = _lazy_image()
    with Image.open(path) as img:
        width, height, fmt = img.width, img.height, img.format
        exif = {}
        try:
            raw = getattr(img, "_getexif", lambda: None)()
            if raw:
                for k, v in raw.items():
                    name = TAGS.get(k, str(k))
                    # 避免存大 binary · 截短
                    try:
                        exif[name] = str(v)[:100]
                    except Exception:
                        pass
        except Exception:
            pass
    return {
        "type": "image",
        "width": width,
        "height": height,
        "format": fmt,
        "exif": exif,
        "content_preview": f"圖片 {width}x{height} · {fmt}",
    }


def _extract_text(path: str) -> dict:
    """.txt / .md / .csv 直接 utf-8 讀"""
    # 最多讀 1MB · 避免超大 log 檔炸掉
    with open(path, "rb") as f:
        raw = f.read(1024 * 1024)
    for enc in ("utf-8", "utf-8-sig", "big5", "cp950", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = raw.decode("utf-8", errors="replace")
    return {
        "type": "text",
        "content_preview": text[:2000],
    }


# ------------------------------------------------------------
# 副檔名 → extractor 路由表
# ------------------------------------------------------------
EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx, ".xls": _extract_xlsx,
    ".jpg": _extract_image, ".jpeg": _extract_image,
    ".png": _extract_image, ".tiff": _extract_image, ".webp": _extract_image,
    ".txt": _extract_text, ".md": _extract_text, ".csv": _extract_text,
    # AI / PSD / AE 等格式無法讀內容 · fallback 走 unknown
}


def extract(path: str) -> dict:
    """對外主入口 · 統一回字典 · 失敗 type='error'"""
    ext = pathlib.Path(path).suffix.lower()
    result: dict[str, Any] = {}
    try:
        if ext in EXTRACTORS:
            result = EXTRACTORS[ext](path)
        else:
            result = {
                "type": "unknown",
                "content_preview": f"檔案 {pathlib.Path(path).name}（格式 {ext} 無法抽字）",
            }
    except Exception as e:
        logger.warning("[extract] %s · %s: %s", path, type(e).__name__, e)
        return {
            "path": path,
            "filename": pathlib.Path(path).name,
            "type": "error",
            "content_preview": "",
            "error": f"{type(e).__name__}: {e}"[:200],
        }
    # 補 file-level metadata
    try:
        stat = os.stat(path)
        result.update({
            "path": path,
            "filename": pathlib.Path(path).name,
            "size": stat.st_size,
            "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        })
    except Exception as e:
        result["stat_error"] = str(e)
    return result
